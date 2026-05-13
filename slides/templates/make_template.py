#!/usr/bin/env python3
"""
Generate a MongoDB-branded reference.pptx for pandoc --reference-doc.

Brand palette (matches mongodb-revealjs.css):
  Slate        #001E2B  — background
  Spring Green #00ED64  — headings, accents, bullets
  Forest Green #00684A  — secondary accents, table headers
  White        #FFFFFF  — body text
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette ──────────────────────────────────────────────────────────
SLATE        = RGBColor(0x00, 0x1E, 0x2B)
SPRING_GREEN = RGBColor(0x00, 0xED, 0x64)
FOREST_GREEN = RGBColor(0x00, 0x68, 0x4A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xCC, 0xCC)

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

HEADING_FONT  = "Calibri"   # closest built-in to Source Serif 4
BODY_FONT     = "Calibri"   # closest built-in to Lexend Deca


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def set_solid_fill(element, rgb: RGBColor):
    """Set solid fill on any DrawingML shape/background element."""
    spPr = element.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(element, qn("p:spPr"))
    # Remove existing fill
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill"):
        for old in spPr.findall(tag):
            spPr.remove(old)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", rgb_to_hex(rgb))


def set_slide_background(slide, rgb: RGBColor):
    """Set the solid background colour of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, left, top, width, height, fill_rgb: RGBColor, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()  # no line
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_name, font_size, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def build_title_slide(prs):
    """Title slide: full slate background, centred green title."""
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    set_slide_background(slide, SLATE)

    # Green accent bar at bottom
    add_rect(slide,
             left=0, top=SLIDE_H - Inches(0.18),
             width=SLIDE_W, height=Inches(0.18),
             fill_rgb=SPRING_GREEN)

    # Green accent bar at top
    add_rect(slide,
             left=0, top=0,
             width=SLIDE_W, height=Inches(0.08),
             fill_rgb=SPRING_GREEN)

    # Title
    add_textbox(slide,
                left=Inches(1.2), top=Inches(2.4),
                width=Inches(10.9), height=Inches(1.5),
                text="Presentation Title",
                font_name=HEADING_FONT, font_size=44, bold=True,
                color=SPRING_GREEN, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide,
                left=Inches(1.2), top=Inches(4.0),
                width=Inches(10.9), height=Inches(0.8),
                text="Subtitle · Author · Date",
                font_name=BODY_FONT, font_size=22, bold=False,
                color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    return slide


def build_content_slide(prs, title_text="Slide Title", body_text="• Bullet point\n• Another point"):
    """Standard content slide: slate background, green title, white body."""
    layout = prs.slide_layouts[1]
    slide  = prs.slides.add_slide(layout)
    set_slide_background(slide, SLATE)

    # Top accent strip
    add_rect(slide,
             left=0, top=0,
             width=SLIDE_W, height=Inches(0.08),
             fill_rgb=SPRING_GREEN)

    # Bottom accent strip
    add_rect(slide,
             left=0, top=SLIDE_H - Inches(0.12),
             width=SLIDE_W, height=Inches(0.12),
             fill_rgb=FOREST_GREEN)

    # Title
    add_textbox(slide,
                left=Inches(0.6), top=Inches(0.25),
                width=Inches(12.0), height=Inches(0.9),
                text=title_text,
                font_name=HEADING_FONT, font_size=32, bold=True,
                color=SPRING_GREEN)

    # Separator line below title
    add_rect(slide,
             left=Inches(0.6), top=Inches(1.15),
             width=Inches(12.0), height=Emu(36000),  # ~0.04 in
             fill_rgb=FOREST_GREEN)

    # Body text
    add_textbox(slide,
                left=Inches(0.6), top=Inches(1.3),
                width=Inches(12.0), height=Inches(5.8),
                text=body_text,
                font_name=BODY_FONT, font_size=22, bold=False,
                color=WHITE)

    return slide


def build_section_slide(prs, section_text="Section"):
    """Full-bleed green slide for section breaks."""
    layout = prs.slide_layouts[1]
    slide  = prs.slides.add_slide(layout)
    set_slide_background(slide, FOREST_GREEN)

    # Bright green bar
    add_rect(slide,
             left=0, top=SLIDE_H - Inches(0.25),
             width=SLIDE_W, height=Inches(0.25),
             fill_rgb=SPRING_GREEN)

    add_textbox(slide,
                left=Inches(1.0), top=Inches(2.8),
                width=Inches(11.3), height=Inches(1.5),
                text=section_text,
                font_name=HEADING_FONT, font_size=40, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    return slide


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_content_slide(prs,
                        title_text="Content Slide",
                        body_text=(
                            "• Main bullet point\n"
                            "• Supporting detail\n"
                            "• Another point\n"
                            "• Keep to 4–6 bullets per slide"
                        ))
    build_content_slide(prs,
                        title_text="Two-Column or Code Slide",
                        body_text="Code block or two-column layout here")
    build_section_slide(prs, section_text="Section Break")

    out = "reference.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
